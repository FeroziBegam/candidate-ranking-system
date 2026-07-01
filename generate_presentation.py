"""
Presentation Generator
Creates PowerPoint deck explaining the approach and results
Can be saved as PDF using PowerPoint or LibreOffice
"""

import json
from pathlib import Path
from datetime import datetime

def generate_presentation(results_json: str, output_pptx: str):
    """Generate PowerPoint presentation from results"""
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
    except ImportError:
        print("python-pptx not installed")
        print("  Install with: pip install python-pptx")
        return False
    
    # Load results
    try:
        with open(results_json, 'r', encoding='utf-8') as f:
            data = json.load(f)
    except Exception as e:
        print(f"Error loading results JSON: {e}")
        return False
    
    # Safely get the results pool
    results = data.get('results', [])[:50]  # Top 50 for slide showcase
    job_req = data.get('job_requirements', {})
    weights = data.get('ranking_weights', {})
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    def add_title_slide(title, subtitle):
        """Add title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(20, 40, 80)  # Dark blue
        
        # Title text block
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle text block
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(200, 220, 255)
        p2.alignment = PP_ALIGN.CENTER
        return slide

    def add_content_slide(title, points):
        """Add standard bullet content slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title text frame
        txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(8.5), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(20, 40, 80)
        
        # Body text frame
        bodyBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(8.5), Inches(5))
        body_tf = bodyBox.text_frame
        body_tf.word_wrap = True
        
        for i, pt_text in enumerate(points):
            if i == 0:
                p_body = body_tf.paragraphs[0]
            else:
                p_body = body_tf.add_paragraph()
            p_body.text = pt_text
            p_body.font.size = Pt(18)
            p_body.space_after = Pt(14)
            p_body.font.color.rgb = RGBColor(50, 50, 50)
        return slide

   # SLIDE 1: TITLE
    add_title_slide(
        "Intelligent Candidate Discovery & Ranking",
        f"A 4-Dimensional AI Recruiting Engine\nGenerated on {datetime.now().strftime('%Y-%m-%d')}"
    )
    
      # SLIDE 2: THE CHALLENGE
    add_content_slide(
        "The Challenge: Beyond Keyword Matching",
        [
            "Recruiters sort through hundreds of profiles and still miss optimal fits.",
            "Standard boolean/keyword filters fail to see what actually matters.",
            "Trap 1: Keyword Stuffers - Profiles listing spammed AI keywords without actual experience.",
            "Trap 2: Plain-Language Engineers - High-quality builders who explain systems in plain text without spammed buzzwords.",
            "Trap 3: Behavioral/Honeypot Profiles - Profiles that look perfect on paper but are stale, inactive, or systematically falsified."
        ]
    )
    
     # SLIDE 3: OUR SOLUTION ARCHITECTURE
    add_content_slide(
        "Our Solution: 4-Dimensional Alignment Engine",
        [
            "Designed to mimic the evaluation process of an expert senior recruiter.",
            f"1. Technical Skill Relevance Alignment ({int(weights.get('skill', 0.40)*100)}% Weight)",
            f"2. Career Trajectory & Experience Fit ({int(weights.get('experience', 0.30)*100)}% Weight)",
            f"3. Operational Momentum & Growth Velocity ({int(weights.get('growth', 0.15)*100)}% Weight)",
            f"4. Platform Behavior & Availability Modifiers ({int(weights.get('culture', 0.15)*100)}% Weight)"
        ]
    )

   # SLIDE 4: THE TOP RECOMMENDATIONS DETAILED SHORTLIST
    shortlist_points = ["Top evaluated profiles extracted by the orchestrator engine pipeline:"]
    for item in results[:5]:
        c_id = item.get('candidate_id', 'Unknown')
        score = item.get('score', 0.0)
        reason = item.get('reasoning', 'No explicit description provided.')
        shortlist_points.append(f"Rank {item.get('rank', '-')}: {c_id} (Score: {score}) - {reason}")
        
    add_content_slide("Top Strategic Recommendation Shortlist", shortlist_points)

   # SLIDE 5: NEXT STEPS & INTERVIEW STRATEGY
    add_content_slide(
        "Strategic Hiring Next Steps",
        [
            "1. Automated Outreach: Target top 10 candidates instantly while platform response indicators are peak.",
            "2. Technical Verification Validation: Deep-dive review of historical product shipping text signals.",
            "3. Reference Integrity Attestation: Match listed career history metrics against standard timeline patterns.",
            "4. Dynamic Optimization: Tune matching weights as hiring data yields active user response rates."
        ]
    )
    
    # SLIDE 6: CONCLUSION
    add_title_slide(
        "Questions?",
        "This ranking system brings clean, verifiable alignment to recruitment.\n Better decisions. Zero keyword noise. High efficiency."
    )
    
    # Save Presentation output
    try:
        prs.save(output_pptx)
        print(f"Presentation saved successfully to: {output_pptx}")
        print("\nTo convert to production ready PDF:")
        print("  Option 1: Open in PowerPoint -> File > Export as PDF")
        print("  Option 2: libreoffice --headless --convert-to pdf output.pptx")
        return True
    except Exception as e:
        print(f"Error saving presentation file: {e}")
        return False

def generate_presentation_from_output(output_dir: str = "output"):
    """Generate presentation from output directory"""
    output_dir = Path(output_dir)
    results_json = output_dir / "ranked_candidates.json"
    output_pptx = output_dir / "presentation.pptx"
    
    if not results_json.exists():
        print(f"Results file not found: {results_json}")
        return False
    
    print("\nGenerating presentation summary deck...")
    return generate_presentation(str(results_json), str(output_pptx))

if __name__ == "__main__":
    import sys
    target_dir = "output"
    if len(sys.argv) > 1:
        target_dir = sys.argv[1]
    generate_presentation_from_output(target_dir)